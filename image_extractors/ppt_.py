from io import BytesIO

from pptx import Presentation
import streamlit as st

from models.llm import LLM_
from utils.utils import encode_image_data


class PPTImageExtractor:
    def __init__(self, llm: LLM_):
        self.llm = llm

    def extract_images_from_pptx(self, file_bytes):
        """
        Extracts images from each slide of a PowerPoint file and generates descriptions.

        Args:
            file_bytes (bytes): The content of the PowerPoint (.pptx) file in bytes.

        Returns:
            dict: A dictionary with slide numbers as keys and a list of (image stream, description) tuples as values.
        """
        images_by_slide = {}
        presentation = Presentation(BytesIO(file_bytes))

        for slide_index, slide in enumerate(presentation.slides):
            slide_images = []
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    try:
                        image_blob = shape.image.blob
                        image_stream = BytesIO(image_blob)
                        base64_image = encode_image_data(image_stream)
                        image_desc = self.llm.ask(base64_image)
                        image_stream.seek(0)
                        slide_images.append((image_stream, image_desc))
                    except Exception as e:
                        st.warning(
                            f"Failed to extract an image on slide {slide_index + 1}: {e}"
                        )
                        raise e
            images_by_slide[slide_index + 1] = slide_images
        return images_by_slide
